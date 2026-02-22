"""Pitch deck generator — 3-stage pipeline: Claude → JSON slides → Jinja2 HTML → PPTX."""

import json
import logging
import os
import re
from html.parser import HTMLParser

import anthropic
from jinja2 import Environment, FileSystemLoader

from backend.config import settings
from backend.matching.pipeline import _build_lead_profile

logger = logging.getLogger(__name__)

_aclient: anthropic.AsyncAnthropic | None = None

# Paths
_TEMPLATE_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "templates")
_OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "generated", "pitchdecks")


def _get_claude_client() -> anthropic.AsyncAnthropic:
    global _aclient
    if _aclient is None:
        _aclient = anthropic.AsyncAnthropic(api_key=settings.anthropic_api_key)
    return _aclient


# ─── Prompt (inline fallback) ────────────────────────────────────────────

_FALLBACK_SYSTEM_PROMPT = """\
You are a sales presentation designer. Given a company profile, product, and match reasoning, \
create a 7-slide pitch deck as JSON.

Return ONLY valid JSON:
{
  "slides": [
    {"slide_number": 1, "title": "...", "body_html": "<p>...</p>", "speaker_notes": "..."}
  ]
}

Slides: 1=Title, 2=Company snapshot, 3=The challenge, 4=The solution, 5=Why this fits, \
6=Social proof, 7=Next steps. Use <h3>, <p>, <ul>/<li>, <strong> in body_html.
"""


def _try_import_prompt() -> str | None:
    """Try to import Person B's prompt from prompts/pitch_deck_prompt.py."""
    try:
        from prompts.pitch_deck_prompt import build_prompt  # type: ignore[import-not-found]

        return build_prompt()
    except (ImportError, ModuleNotFoundError):
        return None


def _parse_json_response(text: str) -> dict:
    """Strip markdown fences and parse JSON from Claude response."""
    try:
        return json.loads(text)  # type: ignore[no-any-return]
    except json.JSONDecodeError:
        pass

    match = re.search(r"```(?:json)?\s*([\s\S]*?)```", text)
    if match:
        return json.loads(match.group(1).strip())  # type: ignore[no-any-return]

    raise ValueError(f"Could not parse JSON from Claude response: {text[:200]}")


def _build_product_summary(product) -> str:
    """Build text summary of a product for the prompt."""
    parts = [f"Product: {product.name}", f"Description: {product.description}"]
    if product.features:
        parts.append(f"Features: {', '.join(product.features)}")
    if product.differentiator:
        parts.append(f"Differentiator: {product.differentiator}")
    if product.example_clients:
        parts.append(f"Example clients: {', '.join(product.example_clients)}")
    if product.current_clients:
        client_names = [c["name"] if isinstance(c, dict) else str(c) for c in product.current_clients]
        parts.append(f"Current clients: {', '.join(client_names)}")
    if product.pricing_model:
        parts.append(f"Pricing: {product.pricing_model}")
    if product.company_name:
        parts.append(f"Sold by: {product.company_name}")
    return "\n".join(parts)


# ─── Stage 1: Claude → JSON slides ──────────────────────────────────────

async def _generate_slides_json(lead, product, match_reasoning: str) -> list[dict]:
    """Generate 7 slides as JSON using Claude."""
    system_prompt = _try_import_prompt() or _FALLBACK_SYSTEM_PROMPT

    lead_profile = _build_lead_profile(lead)
    product_summary = _build_product_summary(product)

    message = await _get_claude_client().messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=3000,
        system=system_prompt,
        messages=[{
            "role": "user",
            "content": (
                f"## Target Company\n{lead_profile}\n\n"
                f"## Product\n{product_summary}\n\n"
                f"## Match Reasoning\n{match_reasoning}\n\n"
                f"Generate a 7-slide pitch deck for selling this product to this company."
            ),
        }],
    )

    response_text: str = message.content[0].text  # type: ignore[union-attr]
    parsed = _parse_json_response(response_text)
    return parsed.get("slides", [])  # type: ignore[no-any-return]


# ─── Stage 2: Jinja2 → HTML ─────────────────────────────────────────────

def _render_html(slides: list[dict], company_name: str, product_name: str) -> str:
    """Render slides into an HTML document using the Jinja2 template."""
    env = Environment(loader=FileSystemLoader(_TEMPLATE_DIR), autoescape=False)
    template = env.get_template("pitch_deck.html")
    return template.render(slides=slides, company_name=company_name, product_name=product_name)


# ─── Stage 3: python-pptx → PPTX ────────────────────────────────────────

# Color constants (matching HTML template dark theme)
_BG_DARK = "0f172a"
_BG_GRADIENT = "1e293b"
_ACCENT = "3b82f6"
_TEXT_LIGHT = "f8fafc"
_TEXT_BODY = "cbd5e1"
_TEXT_H3 = "93c5fd"
_TEXT_BOLD = "60a5fa"
_TEXT_MUTED = "64748b"


def _rgb(hex_color: str):
    """Convert hex string to pptx RGBColor."""
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(hex_color)


class _SlideHTMLParser(HTMLParser):
    """Parse body_html into a list of (text, style) segments for python-pptx."""

    def __init__(self) -> None:
        super().__init__()
        self.segments: list[dict] = []  # {text, bold, color, size, bullet, newpara}
        self._tag_stack: list[str] = []
        self._in_li = False

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        self._tag_stack.append(tag)
        if tag == "p":
            self.segments.append({"text": "", "newpara": True, "size": 18, "color": _TEXT_BODY})
        elif tag == "h3":
            self.segments.append(
                {"text": "", "newpara": True, "size": 22, "color": _TEXT_H3, "bold": True}
            )
        elif tag == "li":
            self._in_li = True
            self.segments.append(
                {"text": "\u25B6  ", "newpara": True, "size": 18, "color": _ACCENT, "bullet": True}
            )
        elif tag == "strong":
            pass  # handled in handle_data

    def handle_endtag(self, tag: str) -> None:
        if self._tag_stack and self._tag_stack[-1] == tag:
            self._tag_stack.pop()
        if tag == "li":
            self._in_li = False

    def handle_data(self, data: str) -> None:
        text = data.strip()
        if not text:
            return
        in_strong = "strong" in self._tag_stack
        in_h3 = "h3" in self._tag_stack

        if self._in_li and self.segments and self.segments[-1].get("bullet"):
            # Append li text as a separate run (body color, not accent)
            self.segments.append(
                {
                    "text": text,
                    "size": 18,
                    "color": _TEXT_BOLD if in_strong else _TEXT_BODY,
                    "bold": in_strong,
                }
            )
        elif in_strong:
            self.segments.append({"text": text, "size": 18, "color": _TEXT_BOLD, "bold": True})
        elif in_h3:
            if self.segments and self.segments[-1].get("color") == _TEXT_H3:
                self.segments[-1]["text"] += text
            else:
                self.segments.append(
                    {"text": text, "newpara": True, "size": 22, "color": _TEXT_H3, "bold": True}
                )
        else:
            self.segments.append({"text": text, "size": 18, "color": _TEXT_BODY})


def _parse_body_html(body_html: str) -> list[dict]:
    """Parse HTML body into styled segments."""
    parser = _SlideHTMLParser()
    parser.feed(body_html)
    return parser.segments


def _set_slide_bg(slide) -> None:
    """Set solid dark background on a slide."""

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(_BG_DARK)
    # Add subtle gradient effect via a second darker rectangle behind everything
    from pptx.util import Emu

    gradient_shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(0),
        Emu(0),
        slide.slide_layout.slide_master.slide_width,
        slide.slide_layout.slide_master.slide_height,
    )
    gradient_shape.fill.solid()
    gradient_shape.fill.fore_color.rgb = _rgb(_BG_GRADIENT)
    gradient_shape.fill.fore_color.brightness = -0.1
    gradient_shape.line.fill.background()
    # Move to back by setting its XML before all other shapes
    sp_tree = slide.shapes._spTree  # type: ignore[attr-defined]
    sp = gradient_shape._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)  # After background but before other shapes


def _add_accent_bar(slide, left, top, width, height) -> None:
    """Add a thin blue accent rectangle."""

    bar = slide.shapes.add_shape(1, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(_ACCENT)
    bar.line.fill.background()


def _add_footer(slide, slide_num: int, total: int) -> None:
    """Add slide number and branding footer."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    # Slide number — bottom right
    num_box = slide.shapes.add_textbox(Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.4))
    num_frame = num_box.text_frame
    num_p = num_frame.paragraphs[0]
    num_p.text = f"{slide_num} / {total}"
    num_p.font.size = Pt(12)
    num_p.font.color.rgb = _rgb(_TEXT_MUTED)
    num_p.alignment = PP_ALIGN.RIGHT

    # Branding — bottom left
    brand_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(4.0), Inches(0.4))
    brand_frame = brand_box.text_frame
    brand_p = brand_frame.paragraphs[0]
    brand_p.text = "Generated by Stick AI"
    brand_p.font.size = Pt(12)
    brand_p.font.color.rgb = _rgb(_TEXT_MUTED)
    brand_p.alignment = PP_ALIGN.LEFT


def _apply_segments_to_frame(text_frame, segments: list[dict]) -> None:
    """Write parsed HTML segments into a pptx text frame with formatting."""
    from pptx.util import Pt

    first_para = True
    current_para = text_frame.paragraphs[0]

    for seg in segments:
        if seg.get("newpara") and not first_para:
            current_para = text_frame.add_paragraph()

        if seg.get("newpara"):
            first_para = False

        run = current_para.add_run()
        run.text = seg.get("text", "")
        run.font.size = Pt(seg.get("size", 18))
        run.font.color.rgb = _rgb(seg.get("color", _TEXT_BODY))
        if seg.get("bold"):
            run.font.bold = True


def _build_title_slide(slide, title: str, body_html: str, total: int) -> None:
    """Build a centered title slide (slide 1)."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    _set_slide_bg(slide)

    # Centered title
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = _rgb(_TEXT_LIGHT)
    title_p.alignment = PP_ALIGN.CENTER

    # Accent bar under title
    _add_accent_bar(slide, Inches(5.5), Inches(3.8), Inches(2.3), Inches(0.06))

    # Subtitle / body
    body_text = re.sub(r"<[^>]+>", "", body_html).strip()
    if body_text:
        sub_box = slide.shapes.add_textbox(Inches(2.0), Inches(4.1), Inches(9.3), Inches(2.0))
        sub_frame = sub_box.text_frame
        sub_frame.word_wrap = True
        sub_p = sub_frame.paragraphs[0]
        sub_p.text = body_text
        sub_p.font.size = Pt(22)
        sub_p.font.color.rgb = _rgb(_TEXT_MUTED)
        sub_p.alignment = PP_ALIGN.CENTER

    # Bottom accent bar
    _add_accent_bar(slide, Inches(0), Inches(7.3), Inches(13.333), Inches(0.06))
    _add_footer(slide, 1, total)


def _build_content_slide(slide, title: str, body_html: str, slide_num: int, total: int) -> None:
    """Build a content slide with left accent bar and parsed HTML body."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    _set_slide_bg(slide)

    # Vertical accent bar (left of title, matching CSS border-left)
    _add_accent_bar(slide, Inches(0.6), Inches(0.5), Inches(0.06), Inches(1.0))

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.5), Inches(11.5), Inches(1.0))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = _rgb(_TEXT_LIGHT)
    title_p.alignment = PP_ALIGN.LEFT

    # Body — parsed HTML with formatting
    body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11.5), Inches(5.0))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True

    segments = _parse_body_html(body_html)
    if segments:
        _apply_segments_to_frame(body_frame, segments)
    else:
        # Fallback: strip HTML and render plain
        body_text = re.sub(r"<[^>]+>", "", body_html).strip()
        body_p = body_frame.paragraphs[0]
        body_p.text = body_text
        body_p.font.size = Pt(18)
        body_p.font.color.rgb = _rgb(_TEXT_BODY)

    # Bottom accent bar
    _add_accent_bar(slide, Inches(0), Inches(7.3), Inches(13.333), Inches(0.06))
    _add_footer(slide, slide_num, total)


def _generate_pptx(slides: list[dict], lead_id: int, product_id: int) -> str:
    """Generate a professionally styled PPTX file from slides data. Returns file path."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    total = len(slides)

    for i, slide_data in enumerate(slides):
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide_data.get("title", "")
        body_html = slide_data.get("body_html", "")

        if i == 0:
            _build_title_slide(slide, title, body_html, total)
        else:
            _build_content_slide(slide, title, body_html, i + 1, total)

        # Speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data.get("speaker_notes", "")

    os.makedirs(_OUTPUT_DIR, exist_ok=True)
    file_path = os.path.join(_OUTPUT_DIR, f"{lead_id}_{product_id}.pptx")
    prs.save(file_path)
    return file_path


# ─── Public API ──────────────────────────────────────────────────────────

async def generate_pitch_deck(lead, product, match_reasoning: str) -> dict:
    """Full 3-stage pipeline: Claude → HTML → PPTX. Returns slides + pptx_path."""
    # Stage 1: Generate slides JSON
    slides = await _generate_slides_json(lead, product, match_reasoning)

    # Stage 2: Render HTML (for preview)
    _render_html(slides, lead.company_name, product.name)

    # Stage 3: Generate PPTX
    pptx_path = _generate_pptx(slides, lead.id, product.id)

    return {"slides": slides, "pptx_path": pptx_path}
